import os, zipfile
from docx import Document
import xlsxwriter
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PIL import Image, ImageDraw

# 📂 Create folder
os.makedirs("test-files", exist_ok=True)

# 📝 Word: document.docx
doc = Document()
doc.add_heading("Test Document", 0)
doc.add_paragraph("Hello World!\nThis is a dummy Word file for QuickTools testing.")
doc.save("test-files/document.docx")

# 📝 Word Form: form.docx
form = Document()
form.add_heading("Dummy Form", 0)
form.add_paragraph("Name: ____________________")
form.add_paragraph("Email: ____________________")
form.save("test-files/form.docx")

# 📊 Excel: table.xlsx
wb = xlsxwriter.Workbook("test-files/table.xlsx")
ws = wb.add_worksheet()
headers = ["Name", "Age", "Country"]
for col, head in enumerate(headers):
    ws.write(0, col, head)
data = [["Ali", 25, "MY"], ["John", 30, "US"], ["Sara", 28, "UK"]]
for row, vals in enumerate(data, 1):
    for col, val in enumerate(vals):
        ws.write(row, col, val)
wb.close()

# 📑 PowerPoint: slides.pptx
prs = Presentation()
for i in range(1, 4):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = f"Slide {i}"
prs.save("test-files/slides.pptx")

# 📄 Function untuk PDF
def make_pdf(filename, lines):
    c = canvas.Canvas(filename, pagesize=letter)
    width, height = letter
    y = height - 100
    for i, line in enumerate(lines):
        c.drawString(100, y - 20 * i, line)
    c.showPage()
    c.save()

# 🗂 PDF files
make_pdf("test-files/sample-1.pdf", ["Page 1", "Page 2", "Page 3"])
make_pdf("test-files/sample-2.pdf", ["Page 1", "Page 2"])
make_pdf("test-files/sample-3.pdf", [f"Page {i}" for i in range(1, 6)])
make_pdf("test-files/report.pdf", ["This is a dummy report PDF."])
make_pdf("test-files/confidential.pdf", ["Confidential Data"])
make_pdf("test-files/contract.pdf", ["Dummy Contract Agreement"])
make_pdf("test-files/draft.pdf", ["Draft Document Page 1", "Draft Page 2"])
make_pdf("test-files/agreement.pdf", ["Agreement Page 1", "Agreement Page 2"])
make_pdf("test-files/long-report.pdf", [f"Paragraph {i}" for i in range(1, 30)])

# 📷 Dummy scanned PDF (simulate image text)
c = canvas.Canvas("test-files/scanned.pdf", pagesize=letter)
c.setFont("Helvetica-Bold", 20)
c.drawString(150, 500, "Scanned TEXT IMAGE (Dummy)")
c.showPage()
c.save()

# 🖼 Images
for name in ["img1.jpg", "img2.png"]:
    img = Image.new("RGB", (400, 300), color=(200, 200, 200))
    d = ImageDraw.Draw(img)
    d.text((150, 150), name, fill=(0, 0, 0))
    img.save(f"test-files/{name}")

# 📦 Zip all files
with zipfile.ZipFile("test-files.zip", "w") as zf:
    for root, _, files in os.walk("test-files"):
        for file in files:
            zf.write(os.path.join(root, file))
